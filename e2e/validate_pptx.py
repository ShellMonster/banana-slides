#!/usr/bin/env python3
"""
PPTX 文件验证脚本
用于 E2E 测试中验证生成的 PPTX 文件是否有效且包含预期内容
"""
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.exc import PackageNotFoundError

def validate_pptx(file_path: str, min_slides: int = 1, expected_keywords: list = None) -> tuple[bool, str]:
    """
    验证 PPTX 文件
    
    Args:
        file_path: PPTX 文件路径
        min_slides: 最少幻灯片数量
        expected_keywords: 期望包含的关键词列表（可选）
    
    Returns:
        (is_valid, message) 元组
    """
    if not os.path.exists(file_path):
        return False, f"文件不存在: {file_path}"
    
    file_size = os.path.getsize(file_path)
    if file_size < 1000:  # 至少 1KB
        return False, f"文件太小 ({file_size} bytes)，可能已损坏"
    
    try:
        # 尝试打开 PPTX 文件
        prs = Presentation(file_path)
    except PackageNotFoundError:
        return False, "文件不是有效的 PPTX 格式"
    except Exception as e:
        return False, f"无法打开 PPTX 文件: {str(e)}"
    
    # 验证幻灯片数量
    slide_count = len(prs.slides)
    if slide_count < min_slides:
        return False, f"幻灯片数量不足: 期望至少 {min_slides} 页，实际 {slide_count} 页"
    
    # 验证内容（如果提供了关键词）
    if expected_keywords:
        found_keywords = []
        missing_keywords = []
        
        # 检查所有幻灯片的文本内容
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        
        for keyword in expected_keywords:
            if keyword.lower() in all_text.lower():
                found_keywords.append(keyword)
            else:
                missing_keywords.append(keyword)
        
        if missing_keywords:
            return False, f"缺少预期关键词: {', '.join(missing_keywords)}"
    
    # 验证幻灯片是否有内容（至少有一个形状）
    empty_slides = 0
    for i, slide in enumerate(prs.slides):
        if len(slide.shapes) == 0:
            empty_slides += 1
    
    if empty_slides == slide_count:
        return False, "所有幻灯片都是空的（没有内容）"
    
    return True, f"验证通过: {slide_count} 页幻灯片，文件大小 {file_size / 1024:.2f} KB"


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python validate_pptx.py <pptx_file> [min_slides] [keyword1] [keyword2] ...")
        sys.exit(1)
    
    file_path = sys.argv[1]
    min_slides = int(sys.argv[2]) if len(sys.argv) > 2 else 1
    expected_keywords = sys.argv[3:] if len(sys.argv) > 3 else None
    
    is_valid, message = validate_pptx(file_path, min_slides, expected_keywords)
    
    if is_valid:
        print(f"✅ {message}")
        sys.exit(0)
    else:
        print(f"❌ {message}")
        sys.exit(1)

